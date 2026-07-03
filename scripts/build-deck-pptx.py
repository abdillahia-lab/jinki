#!/usr/bin/env python3
"""Assemble the Jinki solution deck into a 16:9 PowerPoint.

Each slide is the pixel-perfect 2x PNG (2560x1440) rendered by scripts/render-deck.mjs,
placed full-bleed on a native 13.333x7.5in slide. The PDF (fonts embedded, selectable)
remains the master; this PPTX is the "drops into any PowerPoint" display artifact and
the HTML in tasks/deck/slides/ stays the editable source.

Usage:
  python3 scripts/build-deck-pptx.py            # light deck  -> jinki-solution-deck.pptx
  python3 scripts/build-deck-pptx.py dark        # dark variant -> jinki-solution-deck-dark.pptx
"""
import glob
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    sys.exit("python-pptx not installed. Run: python3 -m pip install python-pptx")

theme = sys.argv[1] if len(sys.argv) > 1 and sys.argv[1] in ("dark", "light") else None
sfx = f"-{theme}" if theme else ""

ROOT = os.getcwd()
RENDERS = os.path.join(ROOT, "tasks/deck/renders")
OUT = os.path.join(ROOT, f"public/docs/jinki-solution-deck{sfx}.pptx")


def matches(path):
    b = os.path.basename(path)
    if theme:
        return b.endswith(f"-{theme}.png")
    return not (b.endswith("-dark.png") or b.endswith("-light.png"))


pngs = sorted(
    (p for p in glob.glob(os.path.join(RENDERS, "slide-*.png")) if matches(p)),
    key=lambda p: int(re.search(r"slide-(\d+)", os.path.basename(p)).group(1)),
)
if not pngs:
    sys.exit(f"no slide PNGs ({sfx or 'light'}) in {RENDERS} — run scripts/render-deck.mjs first")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # fully blank layout
for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=Inches(13.333), height=Inches(7.5))

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"ASSEMBLED {len(pngs)} slides ({theme or 'light'}) -> {OUT}")
